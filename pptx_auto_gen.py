from pptx import Presentation
from pptx.util import Cm
import sys
from argparse import ArgumentParser
import pandas as pd
import os
from datetime import datetime as dt

path = "https://www.togawa.cs.waseda.ac.jp/uploader/files_2019/"

def gen_pptx(template, dir, zemi_name, zemi_date, presenters):
	if os.path.isfile(template):
		prs = Presentation(template)
	else:
		prs = Presentation()
		slide = prs.slides.add_slide(prs.slide_layouts[1])
		slide.shapes.add_textbox(Cm(2.29), Cm(17.94), Cm(5.15), Cm(1.01))

	page = prs.slides[0]
	zemi_title = page.shapes[0]
	contents = page.placeholders[1]
	footer = page.shapes[2]

	zemi_title.text = zemi_name
	footer.text = zemi_date
	date = dt.strptime(zemi_date, '%Y/%m/%d').strftime('%y%m%d')

	new_dir = dir + "_"+ date
	os.makedirs(new_dir, exist_ok=True)
	os.chdir(new_dir)
	for i, (person, file) in enumerate(presenters):
		run = contents.text_frame.paragraphs[0].add_run()
		if len(presenters) == i + 1 :
			run.text = person
		else:
			run.text = person + "\n"
		run.hyperlink.address = file

	prs.save("{}.pptx".format(date))
	os.chdir("../")

	print("PPTX file is successfully generated.")


def gen_csv(args):
	zemi_list = pd.read_csv(args.schedule)
	print("\nInput data:")
	print(zemi_list)
	print("")
	name_list, warn_cnt = name_reader(args.namelist)
	error_cnt = 0

	print("Generating {} zemi PPTX file(s)...\n".format(len(zemi_list)))

	for i, zemi in zemi_list.iterrows():
		print("Generating {} / {}".format(i+1, len(zemi_list)))

		try:
			date = dt.strptime(zemi["date"], '%Y/%m/%d').strftime('%y%m%d')
		except ValueError:
			print("Error: Invalid format or date: {}".format(zemi["date"]))
			error_cnt += 1
			continue

		people = zemi["presenter":].dropna()
		presenters = []
		for person in people:
			if not person in name_list:
				print("Warning: {} is not in laboratry members list.".format(person))
				warn_cnt += 1
			if "." in person:
				file = person.split('.')[1] + "_" + date + ".pptx"
			else:
				file = person + "_" + date + ".pptx"
			presenters.append((name_list.get(person, person), file))

		print_zemi(zemi["dir"], zemi["name"], zemi["date"], presenters)
		gen_pptx(args.template, zemi["dir"], zemi["name"], zemi["date"], presenters)
		out_bat(args.upload_dir, zemi["dir"]+ "_"+ dt.strptime(zemi["date"], '%Y/%m/%d').strftime('%y%m%d'), presenters)
		print ("-" * 50)

	print("{} zemi PPTX file(s) are generated.".format(len(zemi_list)))
	return (error_cnt, warn_cnt)


def gen_int(args):
	print("Interactive mode.")
	name_list, warn_cnt = name_reader(args.namelist)
	print("Input output directory:")
	out_dir = input().rstrip()
	print("Input zemi name:")
	zemi = input().rstrip()
	error_cnt = 0
	while True:
		print("Input zemi date (YYYY/MM/DD):")
		date = input().rstrip()
		try:
			tmp_date = dt.strptime(date, '%Y/%m/%d').strftime('%y%m%d')
			break
		except ValueError:
			print("Error: Invalid format or date. Input again.")
			error_cnt += 1

	i = 1
	presenters = []

	print("Input presenters. Input only \".\" to finish.")
	while True:
		while True:
			print("Presenter name {}:".format(i))
			tmp_person = input().rstrip()
			if tmp_person == "":
				print("Invalid presenter name. Input again.")
			elif tmp_person == "." and i == 1:
				print("You need to specify 1 presenter at least.")
			elif not tmp_person in name_list and tmp_person != ".":
				print("Warning: \"{}\" is not in laboratry members list.".format(tmp_person))
				warn_cnt += 1
				if input("Continue? (y/n): ").lower() == "y":
					print(name_list.get(tmp_person, tmp_person))
					person = name_list.get(tmp_person, tmp_person)
					break
			else:
				print(name_list.get(tmp_person, tmp_person))
				person = name_list.get(tmp_person, tmp_person)
				break
		if tmp_person == "." and i != 1:
			break

		while True:
			print("PPTX file name \"{}\":".format(i))
			if "." in tmp_person:
				tmp_file = tmp_person.split('.')[1] + "_" + tmp_date + ".pptx"
			else:
				tmp_file = tmp_person + "_" + tmp_date + ".pptx"

			print("If you use {}, press ENTER.".format(tmp_file))
			file = input().rstrip()
			if file == "":
				file = tmp_file
				break
			else:
				break
		if file == ".":
			break;
		i+=1
		presenters.append((person, file))

	print_zemi(out_dir, zemi, date, presenters)
	gen_pptx(args.template, out_dir, zemi, date, presenters)
	return error_cnt, warn_cnt


def name_reader(namelist):
	if os.path.isfile(namelist):
		try:
			name_list = pd.read_csv(namelist).set_index("en")["ja"].to_dict()
			code = 0
		except KeyError :
			print("Warning: Invalid format: name list \"{}\".".format(namelist))
			name_list = {}
			code = 1
	else:
		print("Warning: namelist file \"{}\" not found.".format(namelist))
		name_list = {}
		code = 1
	return (name_list, code)

def print_zemi(out_dir, zemi, date, presenters):
	print("Output dir:\t{}".format(out_dir+"_"+dt.strptime(date, '%Y/%m/%d').strftime('%y%m%d')))
	print("Zemi name:\t{}".format(zemi))
	print("Zemi date:\t{}".format(date))
	for i, p in enumerate(presenters):
		print("Presenter {}:\t{}({})".format(i+1, p[0], p[1]))

def out_bat(upload_dir, out_dir, presenters):
	os.chdir(out_dir)
	f = open("copy.bat","w")
	f.write("@echo off\n")
	f.write("set dest={}\n".format(upload_dir))
	for p in presenters:
		f.write("copy %dest%\\{} {} /Y > nul\n".format(p[1], p[1]))
		# f.write("if errorlevel 1 (\npowershell -Command \"wget {}{} -OutFile {};exit $LASTEXITCODE\" > nul\n)\n".format(path, p[1], p[1]))
		f.write("if errorlevel 1 (\npowershell -Command \"(New-Object -com WScript.Shell).Popup(\\\"Copying {} failed.\\\",0,\\\"Copying failed\\\",48)\" > nul\n) else (\npowershell -Command \"(New-Object -com WScript.Shell).Popup(\\\"Copying {} completed.\\\",0,\\\"Copying completed\\\",64)\" > nul)\n\n".format(p[1], p[1]))
	f.close()
	os.chdir("..")
	print("Specified upload directory: {}".format(upload_dir))

def parser():
	argparser = ArgumentParser(description="Zemi PPTX auto generator")
	argparser.add_argument("-i", "--interactive", dest="interactive", action="store_true", help="use interactive mode")
	argparser.add_argument("-s", "--schedule", dest="schedule", default="schedule.csv", action="store", help="specify zemi schedule file (csv). default: schedule.csv")
	argparser.add_argument("-n", "--namelist", dest="namelist", default="name_dic.csv", action="store", help="specify namelist file (csv). default: name_dic.csv")
	argparser.add_argument("-t", "--template", dest="template", default="template.pptx", action="store", help="specify template file (pptx). default: template.pptx")
	argparser.add_argument("-o", "--out-dir", dest="out_dir", default=".", action="store", help="specify output directory. default: current directory")
	argparser.add_argument("-u", "--upload-dir", dest="upload_dir", default="X:\\2019\\regular_seminar", action="store", help="specify directory containing pptx. default: X:\\2019\\regular_seminar")
	args = argparser.parse_args()
	return args

if __name__ == "__main__":
	args = parser()
	print("Zemi PPTX auto generator")
	warn_tmp = 0

	if not os.path.isfile(args.template):
		warn_tmp = 1
		print("Warning: template file \"{}\" not found. No theme will be used.".format(args.template))

	if args.interactive:
		err, warn = gen_int(args)
	else:
		if not os.path.isfile(args.schedule):
			print("Error: schedule file \"{}\" not found.\nAborted.".format(args.schedule))
			print("To run in interactive mode, run with option -i.")
			print("\nThere are 1 error(s) and {} warning(s).".format(warn_tmp))
			exit()
		err, warn = gen_csv(args)
	print("\nThere are {} error(s) and {} warning(s).".format(err, warn + warn_tmp))